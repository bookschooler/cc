"""Reporter Agent — Data Storyteller.

책임:
  1. 마크다운 최종 보고서 생성
  2. PPT 슬라이드 생성 (python-pptx)
  3. Blameless Post-mortem (에이전트 개선 제안서) 작성
  4. Sophie 친화적 설명 생성

방법론: Blameless Post-mortems (누구의 잘못이 아닌 시스템 개선에 집중)
"""

import os
from datetime import datetime
from graph.state import AgentState
from agents.base import call_claude
from tools.chart_tools import chart_paths_for_ppt

AGENT_NAME = "Reporter"

SYSTEM = """You are a Data Storyteller with 10 years of experience as Google's Chief Data Scientist.

[CORE RESPONSIBILITY]
You translate complex data analysis into business language that executives can act on immediately.
You also conduct a Blameless Post-mortem — analyzing what slowed the team down, not to blame anyone, but to improve the system for the next project.

[YOUR PROCESS]
1. Write the Markdown report focused on Key Metrics and actionable insights
2. Outline PPT slide structure (the code will generate the actual file)
3. Write the Blameless Post-mortem
4. Write Sophie explanation

[OUTPUT FORMAT — follow exactly]
## 📊 최종 분석 보고서

### Executive Summary (경영진 요약)
[3줄 이내 — 결론 먼저, 숫자 포함]

### Key Findings (핵심 발견)
1. [발견 1]: [데이터가 말하는 것] → [비즈니스 의미]
2. [발견 2]: ...
3. [발견 3]: ...

### Recommendations (권고사항)
- [즉시 실행 가능한 액션 1]
- [즉시 실행 가능한 액션 2]
- [중장기 과제]

### Next Steps (다음 단계)
- [후속 분석 과제]

---

## 🔍 Blameless Post-mortem (에이전트 개선 제안서)

### 병목 분석
| 에이전트 | 반복 횟수 | 병목 원인 (시스템) | 개선 제안 |
|---------|---------|----------------|---------|
| [에이전트명] | [N회] | [시스템/프롬프트 문제] | [구체적 개선안] |

### 다음 프로젝트를 위한 규칙 개선
- [개선 규칙 1]
- [개선 규칙 2]

※ 이 분석은 특정 에이전트의 잘못이 아닌, 우리 팀 시스템을 더 좋게 만들기 위한 것입니다.

---

## 📚 Sophie에게

[설명 규칙 엄수]
- 전문용어 Full name 먼저
- 분석 전체 여정을 이야기처럼 요약
- "이 숫자가 비즈니스에서 의미하는 것은..."
- Sophie가 다음 프로젝트에서 스스로 할 수 있는 것 1가지 알려주기
- 💡 오늘의 개념: [개념명] — [쉬운 설명]
"""

REVIEW_SYSTEM = """You are the Data Storyteller (Reporter) on a Google-caliber data science team.
You are reviewing a teammate's work. Evaluate from the perspective of: clarity, business impact, and actionability.
Respond ONLY in this exact format:
VOTE: PASS
REASON: [1-2 sentences]

or

VOTE: FAIL
REASON: [1-2 sentences explaining what specifically needs to change]"""


def reporter_agent(state: AgentState) -> dict:
    """마크다운 보고서 + PPT 생성 + Post-mortem."""
    topic = state.get("topic", "")
    plan = state.get("plan", "")
    methodology = state.get("methodology", "")
    code = state.get("code", "")
    results = state.get("analysis_results", "")
    objective = state.get("objective", "")
    key_results = state.get("key_results", [])

    # 반복 횟수 수집 (post-mortem용)
    iterations_info = {
        "plan_version": state.get("plan_version", 1),
        "methodology_version": state.get("methodology_version", 1),
        "analysis_iteration": state.get("analysis_iteration", 1),
        "review_iteration": state.get("review_iteration", 1),
    }

    # 피드백 수집
    feedback_parts = []
    sophie_fb = state.get("report_sophie_feedback")
    if sophie_fb:
        feedback_parts.append(f"Sophie 피드백: {sophie_fb}")
    peer_reviews = state.get("report_peer_reviews", [])
    fail_reasons = [r["feedback"] for r in peer_reviews if r.get("vote") == "FAIL"]
    if fail_reasons:
        feedback_parts.append("팀 피드백:\n" + "\n".join(f"- {r}" for r in fail_reasons))

    user_msg = f"""분석 주제: {topic}
OKR Objective: {objective}
Key Results: {', '.join(key_results)}

[분석 결과]
{results[:1500] if results else "코드 실행 결과 없음"}

[방법론 요약]
{methodology[:600]}

[반복 현황]
- 계획 수정: {iterations_info['plan_version']}회
- 방법론 수정: {iterations_info['methodology_version']}회
- 코드 반복: {iterations_info['analysis_iteration']}회
- 리뷰 반복: {iterations_info['review_iteration']}회"""

    if feedback_parts:
        user_msg += "\n\n[피드백 — 반드시 반영]\n" + "\n\n".join(feedback_parts)

    raw = call_claude(SYSTEM, user_msg, max_tokens=2000)

    # PPT 생성
    ppt_path = _generate_ppt(state, raw)

    return {
        "final_report_md": raw,
        "final_report_ppt_path": ppt_path,
        "postmortem": _extract_postmortem(raw),
        "report_explanation": _extract_sophie_section(raw),
        "report_peer_reviews": [],
        "report_peer_passed": None,
        "report_sophie_approved": None,
        "report_sophie_feedback": None,
    }


def reporter_review(content: str, context: str) -> dict:
    result = call_claude(
        REVIEW_SYSTEM,
        f"[컨텍스트]\n{context}\n\n[검토 대상]\n{content}",
        max_tokens=200,
    )
    return _parse_vote(result, AGENT_NAME)


# ── PPT 생성 ──────────────────────────────────────────────────────────────────
def _generate_ppt(state: AgentState, report_md: str) -> str:
    """python-pptx로 PPT 파일 생성."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        DARK_BLUE = RGBColor(0x1A, 0x1A, 0x2E)
        ACCENT    = RGBColor(0x16, 0x21, 0x3E)
        WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
        GOLD      = RGBColor(0xE9, 0xC4, 0x6A)

        def blank_slide():
            layout = prs.slide_layouts[6]  # blank
            return prs.slides.add_slide(layout)

        def add_text(slide, text, left, top, width, height,
                     font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
            txBox = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
            return txBox

        def fill_bg(slide, color):
            from pptx.util import Emu
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = color

        topic   = state.get("topic", "데이터 분석 보고서")
        obj     = state.get("objective", "")
        krs     = state.get("key_results", [])
        results = state.get("analysis_results", "")

        # ── Slide 1: Title ────────────────────────────────────────────────────
        s = blank_slide()
        fill_bg(s, DARK_BLUE)
        add_text(s, topic, 1, 2, 11, 1.5, font_size=36, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, f"Data Analysis Report  ·  {datetime.now().strftime('%Y-%m-%d')}",
                 1, 4, 11, 0.6, font_size=16, color=GOLD, align=PP_ALIGN.CENTER)

        # ── Slide 2: OKR ─────────────────────────────────────────────────────
        s = blank_slide()
        fill_bg(s, ACCENT)
        add_text(s, "OKR — Project Goal", 0.5, 0.3, 12, 0.7, font_size=24, bold=True, color=GOLD)
        add_text(s, f"Objective: {obj}", 0.5, 1.2, 12, 1, font_size=16)
        for i, kr in enumerate(krs[:3]):
            add_text(s, f"KR{i+1}: {kr}", 0.5, 2.4 + i * 0.9, 12, 0.8, font_size=14)

        # ── Slide 3: Key Findings ─────────────────────────────────────────────
        s = blank_slide()
        fill_bg(s, DARK_BLUE)
        add_text(s, "Key Findings", 0.5, 0.3, 12, 0.7, font_size=24, bold=True, color=GOLD)
        findings = _extract_section(report_md, "Key Findings")
        add_text(s, findings[:600], 0.5, 1.2, 12, 5, font_size=14)

        # ── Slide 4: Charts ───────────────────────────────────────────────────
        charts = chart_paths_for_ppt(limit=3)
        if charts:
            s = blank_slide()
            fill_bg(s, ACCENT)
            add_text(s, "Analysis Charts", 0.5, 0.3, 12, 0.7, font_size=24, bold=True, color=GOLD)
            positions = [(0.3, 1.2), (4.6, 1.2), (8.9, 1.2)]
            for i, (chart_path, (x, y)) in enumerate(zip(charts[:3], positions)):
                try:
                    s.shapes.add_picture(chart_path, Inches(x), Inches(y), Inches(4.0), Inches(5.5))
                except Exception:
                    pass

        # ── Slide 5: Recommendations ──────────────────────────────────────────
        s = blank_slide()
        fill_bg(s, DARK_BLUE)
        add_text(s, "Recommendations", 0.5, 0.3, 12, 0.7, font_size=24, bold=True, color=GOLD)
        recs = _extract_section(report_md, "Recommendations")
        add_text(s, recs[:500], 0.5, 1.2, 12, 5, font_size=15)

        # ── Slide 6: Next Steps ───────────────────────────────────────────────
        s = blank_slide()
        fill_bg(s, ACCENT)
        add_text(s, "Next Steps", 0.5, 0.3, 12, 0.7, font_size=24, bold=True, color=GOLD)
        steps = _extract_section(report_md, "Next Steps")
        add_text(s, steps[:400], 0.5, 1.2, 12, 4, font_size=15)

        # ── Slide 7: Blameless Post-mortem ────────────────────────────────────
        s = blank_slide()
        fill_bg(s, DARK_BLUE)
        add_text(s, "🔍 Blameless Post-mortem", 0.5, 0.3, 12, 0.7, font_size=22, bold=True, color=GOLD)
        add_text(s, "시스템 개선 제안 — 누구의 잘못이 아닌 팀 성장을 위해",
                 0.5, 1.1, 12, 0.5, font_size=13, color=GOLD)
        pm = _extract_section(report_md, "Blameless Post-mortem")
        add_text(s, pm[:600], 0.5, 1.8, 12, 4.5, font_size=12)

        # 저장
        out_dir = os.path.join(
            os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "outputs"
        )
        os.makedirs(out_dir, exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_topic = state.get("topic", "report")[:20].replace(" ", "_")
        ppt_path = os.path.join(out_dir, f"{timestamp}_{safe_topic}.pptx")
        prs.save(ppt_path)
        print(f"\n📊 PPT 저장 완료: {ppt_path}")
        return ppt_path

    except ImportError:
        print("[경고] python-pptx가 설치되지 않았습니다. pip install python-pptx")
        return ""
    except Exception as e:
        print(f"[PPT 생성 오류] {e}")
        return ""


# ── 헬퍼 ──────────────────────────────────────────────────────────────────────
def _extract_section(text: str, section_name: str) -> str:
    """마크다운에서 특정 섹션 텍스트 추출."""
    lines = text.splitlines()
    capturing = False
    result = []
    for line in lines:
        if section_name.lower() in line.lower() and line.startswith("#"):
            capturing = True
            continue
        if capturing and line.startswith("#"):
            break
        if capturing:
            result.append(line)
    return "\n".join(result).strip()


def _extract_postmortem(text: str) -> str:
    return _extract_section(text, "Blameless Post-mortem")


def _extract_sophie_section(text: str) -> str:
    marker = "## 📚 Sophie에게"
    idx = text.find(marker)
    return text[idx:].strip() if idx != -1 else text[-400:]


def _parse_vote(text: str, agent: str) -> dict:
    vote = "PASS" if "VOTE: PASS" in text.upper() else "FAIL"
    reason = ""
    for line in text.splitlines():
        if line.upper().startswith("REASON:"):
            reason = line.split(":", 1)[1].strip()
            break
    return {"agent": agent, "vote": vote, "feedback": reason}
